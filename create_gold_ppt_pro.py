from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

# ==========================================
# กำหนดสี (Corporate Theme: Navy Blue & Gold)
# ==========================================
TITLE_COLOR = RGBColor(0, 51, 102) # Deep Navy Blue
GOLD_COLOR = RGBColor(184, 134, 11) # Dark Goldenrod
TEXT_COLOR = RGBColor(64, 64, 64) # Dark Gray

def style_title(shape, text):
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            run.font.bold = True
            run.font.color.rgb = TITLE_COLOR

def style_content(shape, bullet_points):
    tf = shape.text_frame
    tf.clear() # Clear default formatting
    for bp in bullet_points:
        p = tf.add_paragraph()
        p.text = bp['text']
        p.level = bp['level']
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(22 if bp['level'] == 0 else 18)
            run.font.color.rgb = TITLE_COLOR if bp['level'] == 0 else TEXT_COLOR
            if bp['level'] == 0:
                run.font.bold = True

# ==========================================
# Slide 0: หน้าปก (Title Slide)
# ==========================================
slide_layout_title = prs.slide_layouts[0]
slide0 = prs.slides.add_slide(slide_layout_title)
title0 = slide0.shapes.title
subtitle0 = slide0.placeholders[1]

title0.text = "GOLD INVESTMENT STRATEGY"
for run in title0.text_frame.paragraphs[0].runs:
    run.font.name = 'Arial'
    run.font.color.rgb = GOLD_COLOR
    run.font.bold = True

subtitle0.text = "บทวิเคราะห์และโอกาสการลงทุนทองคำ\nExecutive Summary for Investors"
for run in subtitle0.text_frame.paragraphs[0].runs:
    run.font.name = 'Arial'
    run.font.color.rgb = TITLE_COLOR

# ==========================================
# Slide 1: ปัจจัยผลกระทบ (Key Catalysts)
# ==========================================
slide_layout = prs.slide_layouts[1]
slide1 = prs.slides.add_slide(slide_layout)
style_title(slide1.shapes.title, "1. ปัจจัยขับเคลื่อนราคาทองคำ (Key Catalysts)")
factors = [
    {'text': "นโยบายการเงินและอัตราดอกเบี้ย (Monetary Policy)", 'level': 0},
    {'text': "การปรับลดอัตราดอกเบี้ยของธนาคารกลาง (FED) ส่งผลให้ต้นทุนค่าเสียโอกาสลดลง เป็นบวกต่อทองคำ", 'level': 1},
    {'text': "สภาวะเงินเฟ้อ (Inflationary Pressures)", 'level': 0},
    {'text': "ทองคำยังคงทำหน้าที่เป็นสินทรัพย์ป้องกันความเสี่ยง (Safe Haven) ที่แข็งแกร่ง", 'level': 1},
    {'text': "ความเสี่ยงทางภูมิรัฐศาสตร์ (Geopolitical Risks)", 'level': 0},
    {'text': "ความขัดแย้งระหว่างประเทศกระตุ้นให้เกิดแรงซื้อเพื่อลดความผันผวนของพอร์ตลงทุน", 'level': 1},
    {'text': "การเข้าซื้อของธนาคารกลาง (Central Bank Accumulation)", 'level': 0},
    {'text': "ธนาคารกลางทั่วโลกเพิ่มสัดส่วนทุนสำรองเป็นทองคำอย่างต่อเนื่องเพื่อกระจายความเสี่ยงจาก USD", 'level': 1},
]
style_content(slide1.placeholders[1], factors)

# ==========================================
# Slide 2: แนวโน้ม (Market Outlook)
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
style_title(slide2.shapes.title, "2. ฉากทัศน์และแนวโน้มตลาด (Market Scenarios)")
outlooks = [
    {'text': "📈 Bull Case (กรณีขาขึ้น / เชิงบวก)", 'level': 0},
    {'text': "เศรษฐกิจชะลอตัวเร็วกว่าคาดการณ์ และเข้าสู่วัฏจักรดอกเบี้ยขาลงชัดเจน หนุนเงินทุนไหลเข้า", 'level': 1},
    {'text': "⚖️ Base Case (กรณีฐาน / ทรงตัว)", 'level': 0},
    {'text': "ตลาดรอประเมินทิศทางจากตัวเลขเศรษฐกิจ (CPI, ตลาดแรงงาน) ราคาเคลื่อนไหวในกรอบ (Sideways)", 'level': 1},
    {'text': "📉 Bear Case (กรณีปรับฐาน / เชิงลบ)", 'level': 0},
    {'text': "เศรษฐกิจสหรัฐฯ แข็งแกร่งเกินคาด (No Landing) ดอลลาร์แข็งค่า กดดันราคาทองคำ", 'level': 1},
]
style_content(slide2.placeholders[1], outlooks)

# ==========================================
# Slide 3: ทางเลือกลงทุน (Investment Vehicles)
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
style_title(slide3.shapes.title, "3. เครื่องมือและกลยุทธ์การลงทุน (Investment Vehicles)")
options = [
    {'text': "ทองคำแท่ง (Physical Gold)", 'level': 0},
    {'text': "เหมาะสำหรับการลงทุนระยะยาว ถือครองเป็นสินทรัพย์เพื่อส่งต่อความมั่งคั่ง (Wealth Preservation)", 'level': 1},
    {'text': "กองทุนรวม & ETF (Gold Funds/ETFs)", 'level': 0},
    {'text': "สภาพคล่องสูง บริหารจัดการพอร์ตง่าย ไม่ต้องรับภาระในการจัดเก็บรักษา", 'level': 1},
    {'text': "ตราสารอนุพันธ์ (Gold Futures / Options)", 'level': 0},
    {'text': "เครื่องมือสำหรับป้องกันความเสี่ยง (Hedging) และเพิ่มอัตราทด (Leverage) สำหรับนักลงทุนที่รับความเสี่ยงได้สูง", 'level': 1},
    {'text': "แอปพลิเคชันสะสมทอง (Digital Gold)", 'level': 0},
    {'text': "เหมาะสำหรับกลยุทธ์ถัวเฉลี่ยต้นทุน (DCA) กระจายความเสี่ยงอย่างเป็นระบบ", 'level': 1},
]
style_content(slide3.placeholders[1], options)

prs.save('/home/node/.openclaw/workspace/Gold_Investment_Presentation_Pro.pptx')
print("Professional PowerPoint generated!")