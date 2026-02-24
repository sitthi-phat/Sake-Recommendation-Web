from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# ==========================================
# Slide 1: เหตุการณ์ที่มีผลกระทบ
# ==========================================
slide_layout = prs.slide_layouts[1] # Title and Content
slide1 = prs.slides.add_slide(slide_layout)
title1 = slide1.shapes.title
content1 = slide1.placeholders[1]

title1.text = "ปัจจัยและเหตุการณ์ที่ส่งผลกระทบต่อราคาทองคำ"
tf1 = content1.text_frame
tf1.text = "นโยบายการเงินของธนาคารกลาง (เช่น Fed) - ทิศทางอัตราดอกเบี้ย"

p = tf1.add_paragraph()
p.text = "อัตราเงินเฟ้อ (Inflation) - การใช้ทองคำเป็นสินทรัพย์ป้องกันเงินเฟ้อ"

p = tf1.add_paragraph()
p.text = "ความขัดแย้งทางภูมิรัฐศาสตร์ (Geopolitical Tensions) - สงครามและความไม่สงบ"

p = tf1.add_paragraph()
p.text = "ความเคลื่อนไหวของค่าเงินดอลลาร์สหรัฐ (USD) - มักผกผันกับราคาทองคำ"

p = tf1.add_paragraph()
p.text = "แรงซื้อจากธนาคารกลางทั่วโลก (Central Bank Demand)"

# ==========================================
# Slide 2: แนวโน้มที่จะเป็น
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "แนวโน้มและทิศทางราคาทองคำ"
tf2 = content2.text_frame
tf2.text = "1. ทิศทางขาขึ้น (Bullish Trend)"

p = tf2.add_paragraph()
p.text = "- หากมีการลดอัตราดอกเบี้ย หรือเกิดวิกฤตเศรษฐกิจ/ความไม่สงบ"
p.level = 1

p = tf2.add_paragraph()
p.text = "2. ทิศทางทรงตัว (Sideways)"

p = tf2.add_paragraph()
p.text = "- ตลาดรอประเมินตัวเลขเศรษฐกิจที่สำคัญ เช่น การจ้างงาน, CPI"
p.level = 1

p = tf2.add_paragraph()
p.text = "3. ทิศทางขาลง (Bearish Trend)"

p = tf2.add_paragraph()
p.text = "- หากเศรษฐกิจฟื้นตัวแข็งแกร่ง ดอลลาร์แข็งค่า และนักลงทุนย้ายไปสินทรัพย์เสี่ยง (เช่น หุ้น)"
p.level = 1

# ==========================================
# Slide 3: Option ในการลงทุน
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "ทางเลือกและรูปแบบในการลงทุนทองคำ"
tf3 = content3.text_frame
tf3.text = "1. ทองคำแท่ง / ทองรูปพรรณ (Physical Gold)"

p = tf3.add_paragraph()
p.text = "- เหมาะสำหรับการสะสมระยะยาว จับต้องได้"
p.level = 1

p = tf3.add_paragraph()
p.text = "2. กองทุนรวมทองคำ (Gold Mutual Funds & ETFs)"

p = tf3.add_paragraph()
p.text = "- ใช้เงินเริ่มต้นน้อย ซื้อขายง่ายผ่านแอปธนาคารหรือตลาดหลักทรัพย์"
p.level = 1

p = tf3.add_paragraph()
p.text = "3. แอปพลิเคชันออมทอง (Digital Gold Saving)"

p = tf3.add_paragraph()
p.text = "- ทยอยซื้อสะสมรายเดือน (DCA) ได้ด้วยเงินจำนวนหลักร้อยบาท"
p.level = 1

p = tf3.add_paragraph()
p.text = "4. สัญญาซื้อขายล่วงหน้า (Gold Futures)"

p = tf3.add_paragraph()
p.text = "- เหมาะสำหรับนักลงทุนเก็งกำไร สามารถทำกำไรได้ทั้งขาขึ้นและขาลง (ความเสี่ยงสูง)"
p.level = 1

# Save the presentation
prs.save('/home/node/.openclaw/workspace/Gold_Investment_Plan.pptx')
print("PowerPoint file generated successfully!")
